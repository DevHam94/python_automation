import copy

from pptx import Presentation


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide0 = self.ppt_file.slides[slide_no]

        for shape in slide0.shapes:
            print(f"slide_no:{slide_no} {shape.text}")

    def copy_slide(self, from_slide_no=0, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]

        to_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        for shape in from_slide.shapes:
            el = shape.element
            new_element = copy.deepcopy(el)
            to_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    def duplicate_n_slides(self, slide_cnt, from_slide_no):
        for _ in range(slide_cnt - 1):
            self.copy_slide(from_slide_no=from_slide_no)

    def save(self, filename):
        self.ppt_file.save(filename)

    def change_text(self, slide_no, shape_name, target_text):
        slide = self.ppt_file.slides[slide_no]]
        shape_map = {}
        for shape in slide.shapes:
            print("shape name:", shape.name)

        slide.shapes[0].text = target_text



if __name__ == '__main__':
    ppt_al = PowerPointAutoLabel("재물 조사표.pptx")
    #ppt_al.print_slide_shapes(0)
    ppt_al.duplicate_n_slides(1)
    ppt_al.change_text(0, "", "hello")
    ppt_al.save("[Auto]재물조사표.pptx")
